"""
text height measurement using binary search overflow detection
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from typing import Dict, Any, List, Optional
from pathlib import Path
from PIL import ImageFont
from src.config.poster_config import load_config
from utils.display_formula import (
    CONTENT_BLOCK_GAP,
    FORMULA_VERTICAL_PADDING,
    measure_formula,
    split_content_blocks,
)
from utils.inline_math import (
    apply_script_format,
    split_inline_math,
    visible_inline_math_text,
)

def get_font_file_path(font_name: str) -> str:
    font_mapping = {
        "Arial": "fonts/Arial.ttf",
        "Helvetica Neue": "fonts/HelveticaNeue.ttf",
    }
    
    font_file = font_mapping.get(font_name, "fonts/Arial.ttf")
    project_root = Path(__file__).parent.parent.parent
    font_path = project_root / font_file
    
    return str(font_path)


def estimate_wrapped_line_count(
    text_content: str,
    width_inches: float,
    font_name: str,
    font_size: int,
    margins: Optional[Dict[str, float]] = None,
    font_width_factor: float = 1.0,
    render_width_factor: float = 1.0,
) -> int:
    """Estimate word-wrapped lines with explicit font-width reserves."""

    config = load_config()
    effective_margins = dict(
        config["text_measurement"]["margins"]
    )
    if margins is not None:
        effective_margins.update(margins)

    available_width = max(
        1.0,
        (
            width_inches
            - effective_margins["left"]
            - effective_margins["right"]
        )
        * 72
        * render_width_factor,
    )
    font = ImageFont.truetype(
        get_font_file_path(font_name),
        font_size,
    )

    line_count = 0
    visible_text = "\n".join(
        visible_inline_math_text(line)
        for line in text_content.split("\n")
    )

    for explicit_line in visible_text.split("\n"):
        words = explicit_line.strip().split()
        if not words:
            line_count += 1
            continue

        current_line = ""
        for word in words:
            candidate = (
                f"{current_line} {word}"
                if current_line
                else word
            )
            candidate_width = (
                font.getlength(candidate)
                * font_width_factor
            )
            if current_line and candidate_width > available_width:
                line_count += 1
                current_line = word
            else:
                current_line = candidate

        if current_line:
            line_count += 1

    return max(1, line_count)

def _measure_plain_text_height(
    text_content: str,
    width_inches: float,
    font_name: str = "Arial",
    font_size: int = 44,
    line_spacing: float = 1.0,
    precision: float = 0.001,
    margins: Optional[Dict[str, float]] = None,
) -> Dict[str, Any]:
    """find minimum height for text to fit without font size reduction"""
    
    config = load_config()

    effective_margins = dict(
        config["text_measurement"]["margins"]
    )

    if margins is not None:
        effective_margins.update(margins)

    prs = Presentation()
    slide_layout_index = config["powerpoint"]["slide_layout_blank"]
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
    
    min_height = config["text_measurement"]["min_height"]
    max_height = config["text_measurement"]["max_height"]
    tolerance = precision
    
    while (max_height - min_height) > tolerance:
        test_height = (min_height + max_height) / 2
        
        textbox = slide.shapes.add_textbox(
            left=Inches(config["powerpoint"]["text_frame_positioning"]["default_left"]),
            top=Inches(config["powerpoint"]["text_frame_positioning"]["default_top"]),
            width=Inches(width_inches),
            height=Inches(test_height)
        )
        
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        # Use caller-specific margins when supplied.
        # Body text keeps the configured defaults; section titles
        # use the same compact margins as the renderer.
        text_frame.margin_left = Inches(
            effective_margins["left"]
        )
        text_frame.margin_right = Inches(
            effective_margins["right"]
        )
        text_frame.margin_top = Inches(
            effective_margins["top"]
        )
        text_frame.margin_bottom = Inches(
            effective_margins["bottom"]
        )
        
        # process text exactly like renderer: split by single newlines
        lines = text_content.split('\n')
        
        for line_idx, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            # create paragraph for each line (matching renderer behavior)
            if line_idx == 0 and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = line_spacing

            for segment in split_inline_math(line):
                run = p.add_run()
                run.text = str(segment["text"])
                run.font.name = font_name
                run.font.size = Pt(font_size)
                apply_script_format(
                    run,
                    segment["baseline"],
                    Pt(font_size),
                )
        
        original_size = font_size
        font_reduced = False
        
        try:
            # Use direct font file to bypass cross-platform discovery bug
            font_file_path = get_font_file_path(font_name)
            text_frame.fit_text(font_file=font_file_path, max_size=font_size)
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size.pt < (original_size - 0.5):
                        font_reduced = True
                        break
                if font_reduced:
                    break
        except Exception as e:
            print(f"fit_text error: {e}")
            font_reduced = True
        
        if font_reduced:
            min_height = test_height
        else:
            max_height = test_height
        
        # cleanup textbox
        sp = textbox._element
        sp.getparent().remove(sp)
    
    # calculate newline offset to compensate for pptx rendering discrepancy
    newline_count = text_content.count('\n')
    newline_offset = newline_count * (font_size / 72) * config["text_measurement"]["newline_offset_ratio"]
    final_height = max_height + newline_offset
    
    return {
        "optimal_height": final_height,
        "text_content": text_content,
        "width_inches": width_inches,
        "font_name": font_name,
        "font_size": font_size,
        "line_spacing": line_spacing,
        "precision": precision,
        "newline_count": newline_count,
        "newline_offset": newline_offset
    }


def measure_text_height(
    text_content: str,
    width_inches: float,
    font_name: str = "Arial",
    font_size: int = 44,
    line_spacing: float = 1.0,
    precision: float = 0.001,
    margins: Optional[Dict[str, float]] = None,
) -> Dict[str, Any]:
    """Measure text blocks and vertically stacked display fractions."""

    blocks = split_content_blocks(text_content)
    formula_blocks = [
        block
        for block in blocks
        if block.kind == "formula" and block.formula is not None
    ]
    if not formula_blocks:
        return _measure_plain_text_height(
            text_content=text_content,
            width_inches=width_inches,
            font_name=font_name,
            font_size=font_size,
            line_spacing=line_spacing,
            precision=precision,
            margins=margins,
        )

    config = load_config()
    effective_margins = dict(
        config["text_measurement"]["margins"]
    )
    if margins is not None:
        effective_margins.update(margins)

    formula_width = max(
        0.1,
        width_inches
        - effective_margins["left"]
        - effective_margins["right"],
    )
    total_height = 0.0
    measured_block_count = 0

    for block in blocks:
        block_has_content = (
            (block.kind == "text" and bool(block.text.strip()))
            or (
                block.kind == "formula"
                and block.formula is not None
            )
        )
        if not block_has_content:
            continue
        if measured_block_count:
            total_height += CONTENT_BLOCK_GAP
        measured_block_count += 1

        if block.kind == "text":
            total_height += _measure_plain_text_height(
                text_content=block.text,
                width_inches=width_inches,
                font_name=font_name,
                font_size=font_size,
                line_spacing=line_spacing,
                precision=precision,
                margins=margins,
            )["optimal_height"]
            continue

        formula = block.formula
        if formula is None:
            continue
        if formula.prefix:
            total_height += _measure_plain_text_height(
                text_content=formula.prefix,
                width_inches=width_inches,
                font_name=font_name,
                font_size=font_size,
                line_spacing=line_spacing,
                precision=precision,
                margins=margins,
            )["optimal_height"]

        _, formula_height = measure_formula(
            formula.mathtext,
            font_size,
            formula_width,
        )
        total_height += formula_height + (2 * FORMULA_VERTICAL_PADDING)

    return {
        "optimal_height": total_height,
        "text_content": text_content,
        "width_inches": width_inches,
        "font_name": font_name,
        "font_size": font_size,
        "line_spacing": line_spacing,
        "precision": precision,
        "newline_count": text_content.count("\n"),
        "newline_offset": 0.0,
        "display_formula_count": len(formula_blocks),
    }

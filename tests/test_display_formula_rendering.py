import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from src.agents.renderer import Renderer
from src.layout.text_height_measurement import measure_text_height
from utils.display_formula import (
    extract_display_formula,
    render_formula_png,
)


PRIVACY_LOSS_CLAIM = (
    "The privacy loss for DP-CutMixSL is "
    "ε_CutMix(α) = (α × max_{i∈C} λ_i)/2 × "
    "((Δ^2 × D_s)/σ_s^2 + "
    "(max_{i∈C} λ_i × D_y)/σ_y^2)."
)

SMASHED_TERM_CLAIM = (
    "For DP-CutMixSL, the smashed-data term satisfies "
    "(α/(2 × σ_s^2)) × λ_{i'} × D_s × Δ^2 "
    "≤ (max_{i∈C} λ_i) × Δ^2 × D_s."
)


class DisplayFormulaRenderingTests(unittest.TestCase):
    def test_converts_slashes_to_stacked_fractions_without_reordering(self):
        formula = extract_display_formula(PRIVACY_LOSS_CLAIM)

        self.assertIsNotNone(formula)
        self.assertEqual(
            "The privacy loss for DP-CutMixSL is",
            formula.prefix,
        )
        self.assertEqual(3, formula.mathtext.count(r"\frac"))
        self.assertNotIn("/", formula.mathtext)
        self.assertIn(
            r"\epsilon_{\mathrm{CutMix}}\left(\alpha\right) =",
            formula.mathtext,
        )
        self.assertLess(
            formula.mathtext.index(r"\Delta^{2}"),
            formula.mathtext.index(r"\max_{i\in C}", 80),
        )

    def test_supports_parenthesized_denominator_and_comparison(self):
        formula = extract_display_formula(SMASHED_TERM_CLAIM)

        self.assertIsNotNone(formula)
        self.assertEqual(1, formula.mathtext.count(r"\frac"))
        self.assertIn(
            r"\frac{\alpha}{2 \times \sigma_{s}^{2}}",
            formula.mathtext,
        )
        self.assertIn(r"\leq", formula.mathtext)

    def test_does_not_treat_prose_slashes_as_display_math(self):
        self.assertIsNone(
            extract_display_formula(
                "The privacy/accuracy trade-off is evaluated experimentally."
            )
        )

    def test_measurement_reserves_stacked_fraction_height(self):
        measurement = measure_text_height(
            PRIVACY_LOSS_CLAIM,
            width_inches=16.0,
            font_size=44,
            margins={
                "left": 0.10,
                "right": 0.10,
                "top": 0.05,
                "bottom": 0.05,
            },
        )
        prefix_measurement = measure_text_height(
            "The privacy loss for DP-CutMixSL is",
            width_inches=16.0,
            font_size=44,
        )

        self.assertEqual(1, measurement["display_formula_count"])
        self.assertGreater(
            measurement["optimal_height"],
            prefix_measurement["optimal_height"] + 0.5,
        )

    def test_renderer_embeds_formula_with_exact_source_alt_text(self):
        renderer = Renderer()
        renderer.styling_interfaces = {
            "font_sizes": {"body_text": 44},
            "line_spacing": 1.0,
        }
        content = (
            "A verified lead-in.\n"
            f"{PRIVACY_LOSS_CLAIM}\n"
            "A verified conclusion."
        )
        measurement = measure_text_height(
            content,
            width_inches=16.0,
            font_size=44,
            margins={
                "left": 0.10,
                "right": 0.10,
                "top": 0.05,
                "bottom": 0.05,
            },
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "display-formula.pptx"
            presentation = Presentation()
            presentation.slide_width = Inches(18)
            presentation.slide_height = Inches(10)
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[6]
            )
            renderer._add_enhanced_text(
                slide,
                content,
                Inches(1),
                Inches(1),
                Inches(16),
                Inches(measurement["optimal_height"] + 0.1),
                {
                    "font_family": "Arial",
                    "font_size": 44,
                    "font_color": "#000000",
                    "line_spacing": 1.0,
                },
            )
            presentation.save(output_path)

            reopened = Presentation(output_path)
            pictures = [
                shape
                for shape in reopened.slides[0].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            self.assertEqual(1, len(pictures))
            self.assertEqual(
                extract_display_formula(PRIVACY_LOSS_CLAIM).expression,
                pictures[0]._element.nvPicPr.cNvPr.get("descr"),
            )

    def test_formula_png_contains_rendered_ink(self):
        formula = extract_display_formula(PRIVACY_LOSS_CLAIM)
        png = render_formula_png(
            formula.mathtext,
            44,
            "#000000",
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            image_path = Path(temp_dir) / "formula.png"
            image_path.write_bytes(png)
            image = Image.open(image_path).convert("L")
            self.assertLess(image.getextrema()[0], 32)
            self.assertGreater(image.width, image.height * 4)


if __name__ == "__main__":
    unittest.main()

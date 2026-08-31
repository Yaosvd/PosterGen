import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.agents.renderer import Renderer
from src.layout.text_height_measurement import measure_text_height
from utils.inline_math import (
    SCRIPT_BASELINES,
    SCRIPT_FONT_SCALE,
    split_inline_math,
    visible_inline_math_text,
)


FORMULA = (
    "ε_CutMix(α) = (α × max_{i∈C} λ_i)/2 × "
    "((Δ^2 × D_s)/σ_s^2)."
)


class InlineMathRenderingTests(unittest.TestCase):
    def test_splits_braced_and_unbraced_scripts(self):
        segments = split_inline_math(FORMULA)

        self.assertEqual(
            [
                ("ε", None),
                ("CutMix", "subscript"),
                ("(α) = (α × max", None),
                ("i∈C", "subscript"),
                (" λ", None),
                ("i", "subscript"),
                (")/2 × ((Δ", None),
                ("2", "superscript"),
                (" × D", None),
                ("s", "subscript"),
                (")/σ", None),
                ("s", "subscript"),
                ("2", "superscript"),
                (").", None),
            ],
            [
                (segment["text"], segment["baseline"])
                for segment in segments
            ],
        )

    def test_parentheses_stay_on_the_normal_baseline(self):
        segments = split_inline_math("ε_CutMix(α)")

        self.assertEqual("εCutMix(α)", visible_inline_math_text("ε_CutMix(α)"))
        self.assertEqual("(α)", segments[-1]["text"])
        self.assertIsNone(segments[-1]["baseline"])

    def test_preserves_theorem_one_notation(self):
        theorem = "ε_Mix(α) ≤ ε_CutMix(α) ≤ ε_o(α)"
        segments = split_inline_math(theorem)

        self.assertEqual(
            ["Mix", "CutMix", "o"],
            [
                segment["text"]
                for segment in segments
                if segment["baseline"] == "subscript"
            ],
        )
        self.assertEqual(
            "εMix(α) ≤ εCutMix(α) ≤ εo(α)",
            visible_inline_math_text(theorem),
        )

    def test_escaped_script_markers_stay_literal(self):
        text = r"file\_name and x\^2"

        self.assertEqual(
            "file_name and x^2",
            visible_inline_math_text(text),
        )
        self.assertTrue(all(
            segment["baseline"] is None
            for segment in split_inline_math(text)
        ))

    def test_renderer_persists_editable_powerpoint_baselines(self):
        renderer = Renderer()

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "inline-math.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[6]
            )
            paragraph = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.5),
                Inches(12),
                Inches(1.5),
            ).text_frame.paragraphs[0]
            renderer._parse_and_add_runs(
                paragraph,
                FORMULA,
                "Arial",
                Pt(40),
                RGBColor(0, 0, 0),
            )
            presentation.save(output_path)

            reopened = Presentation(output_path)
            runs = reopened.slides[0].shapes[0].text_frame.paragraphs[0].runs
            subscript_runs = [
                run
                for run in runs
                if run.font._element.get("baseline")
                == str(SCRIPT_BASELINES["subscript"])
            ]
            superscript_runs = [
                run
                for run in runs
                if run.font._element.get("baseline")
                == str(SCRIPT_BASELINES["superscript"])
            ]

            self.assertEqual(
                ["CutMix", "i∈C", "i", "s", "s"],
                [run.text for run in subscript_runs],
            )
            self.assertEqual(
                ["2", "2"],
                [run.text for run in superscript_runs],
            )
            self.assertTrue(all(
                abs(run.font.size.pt - (40 * SCRIPT_FONT_SCALE)) < 0.1
                for run in subscript_runs + superscript_runs
            ))
            self.assertTrue(any(
                "(α)" in run.text
                and run.font._element.get("baseline") is None
                for run in runs
            ))

    def test_math_runs_keep_existing_emphasis(self):
        segments = Renderer()._tokenize_formatting("**ε_CutMix(α)**")

        self.assertEqual(
            [("ε", None), ("CutMix", "subscript"), ("(α)", None)],
            [
                (segment["text"], segment["baseline"])
                for segment in segments
            ],
        )
        self.assertTrue(all(segment["bold"] for segment in segments))

    def test_height_measurement_accepts_inline_math(self):
        measurement = measure_text_height(
            FORMULA,
            width_inches=8.0,
            font_size=44,
        )

        self.assertGreater(measurement["optimal_height"], 0)


if __name__ == "__main__":
    unittest.main()
